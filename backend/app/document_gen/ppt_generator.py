"""
PowerPoint Generator — generates comprehensive bid presentations using python-pptx.
Covers all 17 agent outputs in a structured executive deck.
"""

import os
from typing import Any, Dict, Optional
from datetime import datetime


def _unwrap(output) -> dict:
    """Unwrap agent output envelope {status, agent, result: {...}} if needed."""
    if not output or not isinstance(output, dict):
        return {}
    if "result" in output and isinstance(output.get("result"), dict):
        if "status" in output or "agent" in output:
            return output["result"]
    return output


def _truncate(text: str, max_len: int = 800) -> str:
    """Safely truncate text for slide content."""
    if not text:
        return ""
    text = str(text).replace("**", "").strip()
    if len(text) > max_len:
        return text[: max_len - 3] + "..."
    return text


def _add_bullet_slide(prs, title: str, bullets: list, subtitle: str = ""):
    """Add a standard bullet-point slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = subtitle if subtitle else ""
    for bullet in bullets[:12]:  # Max 12 bullets per slide
        p = tf.add_paragraph()
        p.text = _truncate(str(bullet), 200)
    return slide


def _add_table_slide(prs, title: str, headers: list, rows: list):
    """Add a slide with a data table."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide.shapes.title.text = title

    if not rows:
        return slide

    # Table dimensions
    num_rows = min(len(rows), 10) + 1  # Header + data, max 10 data rows
    num_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(0.4) * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0xFF)

    # Data rows
    for r_idx, row in enumerate(rows[:10]):
        for c_idx, val in enumerate(row[:num_cols]):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = _truncate(str(val) if val else "", 100)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)

    return slide


def generate_proposal_deck(
    bid_data: Dict[str, Any],
    output_dir: str = "../knowledge_base",
    filename: Optional[str] = None,
) -> Dict[str, Any]:
    """Generate a comprehensive bid proposal presentation from all agent outputs."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"status": "error", "message": "python-pptx not installed"}

    prs = Presentation()
    manifest = bid_data.get("manifest", {})
    client = manifest.get("client", {})
    rfp = manifest.get("rfp", {})
    client_name = client.get("name") or bid_data.get("client_name", "Client")
    bidder_name = manifest.get("bidder_profile", {}).get("name", "[Bidder]")
    contract_type = (
        bid_data.get("contract_type", "") or rfp.get("contract_type", "")
    ).upper() or "ENGAGEMENT"
    products = bid_data.get("products", []) or rfp.get("products", [])
    industry = client.get("industry") or bid_data.get("client_industry", "")

    # ═══════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Proposal for {client_name}"
    slide.placeholders[
        1
    ].text = (
        f"{bidder_name}\n{contract_type} Services\n{datetime.now().strftime('%B %Y')}"
    )

    # ═══════════════════════════════════════════════════════
    # SLIDE 2: AGENDA
    # ═══════════════════════════════════════════════════════
    agenda_items = [
        "Executive Summary & Key Requirements",
        "Strategic Assessment & Win Probability",
        "Scope of Work & Effort Estimation",
        "Solution Architecture",
        "Competitive Positioning & Differentiators",
        "Commercial Model & Pricing",
        "Compliance & Risk Assessment",
        "Automation & AI Opportunities",
        "Transition & Change Management",
        "Quality Assurance & Next Steps",
    ]
    _add_bullet_slide(prs, "Agenda", agenda_items)

    # ═══════════════════════════════════════════════════════
    # SLIDE 3: EXECUTIVE SUMMARY
    # ═══════════════════════════════════════════════════════
    intake = _unwrap(bid_data.get("intake_output"))
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    body = slide.placeholders[1]
    tf = body.text_frame

    exec_summary = intake.get("executive_summary", "") or intake.get("narrative", "")
    tf.text = (
        _truncate(exec_summary, 600)
        if exec_summary
        else f"{bidder_name} presents its comprehensive response to the {contract_type} requirements."
    )

    if products:
        p = tf.add_paragraph()
        p.text = f"\nProducts: {', '.join(products)}"
    if industry:
        p = tf.add_paragraph()
        p.text = f"Industry: {industry}"

    # ═══════════════════════════════════════════════════════
    # SLIDE 4: KEY REQUIREMENTS
    # ═══════════════════════════════════════════════════════
    extracted = intake.get("extracted_fields", {})
    if extracted:
        bullets = []
        for k, v in extracted.items():
            if isinstance(v, dict):
                val = v.get("value", "")
                display = ", ".join(val) if isinstance(val, list) else str(val)
                if display:
                    bullets.append(f"{k.replace('_', ' ').title()}: {display}")
        if bullets:
            _add_bullet_slide(prs, "Key Requirements Identified", bullets)

    # ═══════════════════════════════════════════════════════
    # SLIDE 5: STRATEGIC ASSESSMENT
    # ═══════════════════════════════════════════════════════
    bid_score = manifest.get("bid_score", {})
    wp = bid_data.get("win_probability")
    rec = bid_data.get("bid_recommendation", "")

    if bid_score or wp:
        bullets = []
        if rec:
            bullets.append(f"Recommendation: {rec}")
        if wp:
            bullets.append(f"Win Probability: {wp * 100:.0f}%")
        dimensions = bid_score.get("dimensions", {})
        for dim_name, dim_data in dimensions.items():
            if isinstance(dim_data, dict):
                score = dim_data.get("score", "")
                bullets.append(f"{dim_name.replace('_', ' ').title()}: {score}/10")
        narrative = bid_score.get("narrative", "")
        if narrative:
            bullets.append(f"\n{_truncate(narrative, 300)}")
        _add_bullet_slide(prs, "Strategic Assessment", bullets)

    # ═══════════════════════════════════════════════════════
    # SLIDE 6-7: SCOPE OF WORK
    # ═══════════════════════════════════════════════════════
    scope = _unwrap(bid_data.get("scope_output"))
    scope_pkg = scope.get("scope_package", {}) if isinstance(scope, dict) else {}

    scope_narrative = scope.get("narrative", "")
    if scope_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Scope of Work"
        slide.placeholders[1].text = _truncate(scope_narrative, 800)

    # Scope by platform table
    scope_by_platform = scope_pkg.get("scope_by_platform", [])
    if scope_by_platform:
        rows = []
        for plat in scope_by_platform:
            if isinstance(plat, dict):
                rows.append(
                    [
                        plat.get("platform", ""),
                        str(plat.get("platform_effort_days", "")),
                        plat.get("scope_summary", "")[:80],
                    ]
                )
        if rows:
            _add_table_slide(
                prs, "Scope by Platform", ["Platform", "Effort (days)", "Summary"], rows
            )

    # Effort summary
    total_effort = scope_pkg.get("total_effort_days", 0)
    scope_pkg.get("timeline_months", "")
    team = scope_pkg.get("team_model", [])
    if team and isinstance(team, list):
        rows = []
        for r in team:
            if isinstance(r, dict):
                rows.append(
                    [
                        r.get("role", ""),
                        str(r.get("count", "")),
                        r.get("location", ""),
                        r.get("platform", ""),
                    ]
                )
        if rows:
            _add_table_slide(
                prs,
                f"Proposed Team (Total Effort: {total_effort} days)",
                ["Role", "Count", "Location", "Platform"],
                rows,
            )

    # ═══════════════════════════════════════════════════════
    # SLIDE 8: SOLUTION ARCHITECTURE
    # ═══════════════════════════════════════════════════════
    solution = _unwrap(bid_data.get("solution_output"))
    sol_narrative = solution.get("narrative", "")
    if sol_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Solution Architecture"
        slide.placeholders[1].text = _truncate(sol_narrative, 800)

    sol_pkg = solution.get("solution_package", {}) if isinstance(solution, dict) else {}
    overview = sol_pkg.get(
        "solution_overview", solution.get("architecture_overview", "")
    )
    if overview and not sol_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Solution Architecture"
        slide.placeholders[1].text = _truncate(overview, 800)

    # ═══════════════════════════════════════════════════════
    # SLIDE 9: COMPETITIVE POSITIONING
    # ═══════════════════════════════════════════════════════
    competitive = _unwrap(bid_data.get("competitive_output"))
    comp_narrative = competitive.get("narrative", "")
    if comp_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Competitive Positioning"
        slide.placeholders[1].text = _truncate(comp_narrative, 800)

    comp_pkg = (
        competitive.get("competitive_package", {})
        if isinstance(competitive, dict)
        else {}
    )
    win_themes = comp_pkg.get(
        "win_themes",
        competitive.get("win_themes", competitive.get("key_differentiators", [])),
    )
    if win_themes and isinstance(win_themes, list):
        bullets = []
        for wt in win_themes[:8]:
            if isinstance(wt, dict):
                bullets.append(
                    f"✓ {wt.get('theme', wt.get('name', ''))}: {wt.get('description', '')}"
                )
            else:
                bullets.append(f"✓ {wt}")
        if bullets:
            _add_bullet_slide(prs, "Key Differentiators", bullets)

    # ═══════════════════════════════════════════════════════
    # SLIDE 10-11: COMMERCIAL MODEL
    # ═══════════════════════════════════════════════════════
    commercial = _unwrap(bid_data.get("commercial_output"))
    comm_narrative = commercial.get("narrative", "")
    if comm_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Commercial Model"
        slide.placeholders[1].text = _truncate(comm_narrative, 800)

    pl = commercial.get("pl_model", {})
    revenue = pl.get("revenue", {})
    profit = pl.get("profitability", {})
    if revenue or profit:
        bullets = []
        if revenue.get("total_contract_value"):
            tcv = revenue["total_contract_value"]
            bullets.append(
                f"Total Contract Value: ${tcv:,.0f}"
                if isinstance(tcv, (int, float))
                else f"TCV: {tcv}"
            )
        if revenue.get("annual_revenue"):
            ar = revenue["annual_revenue"]
            bullets.append(
                f"Annual Revenue: ${ar:,.0f}"
                if isinstance(ar, (int, float))
                else f"Annual: {ar}"
            )
        if profit.get("margin_percent"):
            bullets.append(f"Gross Margin: {profit['margin_percent']:.1f}%")
        if profit.get("gross_profit"):
            gp = profit["gross_profit"]
            bullets.append(
                f"Gross Profit: ${gp:,.0f}"
                if isinstance(gp, (int, float))
                else f"Profit: {gp}"
            )

        guardrail = commercial.get("margin_guardrail", {})
        if guardrail:
            bullets.append(
                f"\nMargin Guardrail: {guardrail.get('status', 'N/A').upper()} — {guardrail.get('message', '')}"
            )

        if bullets:
            _add_bullet_slide(prs, "Financial Summary", bullets)

    # Resource plan
    resources = commercial.get("resource_plan", [])
    if resources and isinstance(resources, list):
        rows = []
        for r in resources:
            if isinstance(r, dict):
                rows.append(
                    [
                        r.get("role", ""),
                        r.get("location", ""),
                        str(r.get("count", "")),
                        f"M{r.get('start_month', 1)}",
                    ]
                )
        if rows:
            _add_table_slide(
                prs, "Resource Plan", ["Role", "Location", "FTEs", "Start"], rows
            )

    # ═══════════════════════════════════════════════════════
    # SLIDE 12: COMPLIANCE & RISK
    # ═══════════════════════════════════════════════════════
    compliance = _unwrap(bid_data.get("compliance_output"))
    comp_risk_narrative = compliance.get("narrative", "")
    if comp_risk_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Compliance & Risk Assessment"
        slide.placeholders[1].text = _truncate(comp_risk_narrative, 800)

    risk_data = compliance.get("contractual_risks", compliance.get("risk_register", []))
    if risk_data and isinstance(risk_data, list):
        rows = []
        for risk in risk_data[:8]:
            if isinstance(risk, dict):
                rows.append(
                    [
                        risk.get("id", risk.get("clause", "")),
                        risk.get("description", risk.get("risk", ""))[:60],
                        risk.get("severity", risk.get("impact", "")),
                        risk.get("mitigation", "")[:60],
                    ]
                )
        if rows:
            _add_table_slide(
                prs, "Risk Register", ["ID", "Risk", "Severity", "Mitigation"], rows
            )

    # ═══════════════════════════════════════════════════════
    # SLIDE 13: AUTOMATION & AI
    # ═══════════════════════════════════════════════════════
    automation = _unwrap(bid_data.get("automation_ai_output"))
    auto_narrative = automation.get("narrative", "")
    opportunities = automation.get("opportunities", [])
    if auto_narrative or opportunities:
        if auto_narrative:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Automation & AI Opportunities"
            slide.placeholders[1].text = _truncate(auto_narrative, 800)

        if opportunities:
            rows = []
            for opp in opportunities[:8]:
                if isinstance(opp, dict):
                    rows.append(
                        [
                            opp.get("title", opp.get("name", "")),
                            opp.get("platform", ""),
                            opp.get("priority", ""),
                            str(
                                opp.get(
                                    "estimated_savings", opp.get("effort_reduction", "")
                                )
                            ),
                        ]
                    )
            if rows:
                _add_table_slide(
                    prs,
                    "Automation Opportunity Catalogue",
                    ["Opportunity", "Platform", "Priority", "Savings"],
                    rows,
                )

    # ═══════════════════════════════════════════════════════
    # SLIDE 14: TRANSITION & CHANGE
    # ═══════════════════════════════════════════════════════
    transition = _unwrap(bid_data.get("transition_change_output"))
    trans_narrative = transition.get("narrative", "")
    if trans_narrative:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Transition & Change Management"
        slide.placeholders[1].text = _truncate(trans_narrative, 800)

    trans_plan = transition.get("transition_plan", {})
    if isinstance(trans_plan, dict):
        phases = trans_plan.get("phases", [])
        if phases:
            rows = []
            for phase in phases:
                if isinstance(phase, dict):
                    rows.append(
                        [
                            phase.get("phase", phase.get("name", "")),
                            phase.get("duration", ""),
                            phase.get("description", "")[:60],
                        ]
                    )
            if rows:
                _add_table_slide(
                    prs, "Transition Phases", ["Phase", "Duration", "Activities"], rows
                )

    # ═══════════════════════════════════════════════════════
    # SLIDE 15: QA SUMMARY
    # ═══════════════════════════════════════════════════════
    qa = _unwrap(bid_data.get("qa_output"))
    qa_scores = qa.get("quality_scores", {})
    if qa_scores:
        bullets = []
        for k, v in qa_scores.items():
            label = k.replace("_", " ").title()
            score_val = f"{v}/100" if isinstance(v, (int, float)) else str(v)
            bullets.append(f"{label}: {score_val}")
        _add_bullet_slide(prs, "Quality Assurance Scores", bullets)

    # ═══════════════════════════════════════════════════════
    # SLIDE 16: NEXT STEPS
    # ═══════════════════════════════════════════════════════
    next_steps = [
        f"Schedule discovery workshop with {client_name} stakeholders",
        "Validate scope assumptions and technical requirements",
        "Align on governance model and escalation pathways",
        "Finalize commercial terms and contract structure",
        "Commence transition phase per the proposed timeline",
    ]
    _add_bullet_slide(prs, "Recommended Next Steps", next_steps)

    # ═══════════════════════════════════════════════════════
    # SLIDE 17: THANK YOU
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[
        1
    ].text = (
        f"{bidder_name} — Your Strategic Partner\n\n{datetime.now().strftime('%B %Y')}"
    )

    # ═══════════════════════════════════════════════════════
    # SAVE
    # ═══════════════════════════════════════════════════════
    os.makedirs(output_dir, exist_ok=True)
    fname = (
        filename
        or f"Proposal_{client_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
    )
    path = os.path.join(output_dir, fname)
    prs.save(path)

    return {
        "status": "success",
        "path": path,
        "filename": fname,
        "slides": len(prs.slides),
    }
