"""
Knowledge Base Upload & Ingestion Service.
Handles document parsing, text extraction, chunking, embedding, and metadata tagging.
"""

import os
from typing import Any, Dict, List, Optional
from app.knowledge.embeddings import embedding_service
from app.knowledge.rag import rag_pipeline


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".csv",
    ".txt",
    ".md",
    ".pptx",
    ".ppt",
}


def extract_text_from_file(file_path: str) -> str:
    """Extract text content from various file formats."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return _extract_docx(file_path)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(file_path)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(file_path)
    elif ext in (".txt", ".md", ".csv"):
        return _extract_text(file_path)
    else:
        return ""


def extract_text(file_path: str, filename: str = "") -> str:
    """Convenience alias — extracts text using file path (or filename for extension detection)."""
    if filename and not os.path.splitext(file_path)[1]:
        # If file_path has no extension but filename does, use filename extension
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return _extract_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif ext in (".txt", ".md", ".csv"):
            return _extract_text(file_path)
        return ""
    return extract_text_from_file(file_path)


def _extract_pdf(path: str) -> str:
    try:
        import fitz

        doc = fitz.open(path)
        return "\n".join(page.get_text() for page in doc)
    except Exception:
        return ""


def extract_images_from_pdf(path: str, output_dir: str = "") -> list:
    """Extract embedded images from PDF pages.

    Returns a list of dicts: [{page, path, width, height, size_bytes}]
    Filters out images smaller than 10KB (logos, icons, bullets) to focus
    on diagrams, architecture charts, and infographics.
    """
    images = []
    try:
        import fitz

        if not output_dir:
            output_dir = os.path.join(os.path.dirname(path), "extracted_images")
        os.makedirs(output_dir, exist_ok=True)

        doc = fitz.open(path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue
                    image_bytes = base_image["image"]
                    image_ext = base_image.get("ext", "png")
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)

                    # Skip tiny images (logos, icons, bullets)
                    if len(image_bytes) < 10_000:
                        continue
                    # Skip extremely narrow/flat images (likely decorative lines)
                    if width < 100 or height < 100:
                        continue

                    img_filename = f"page{page_num + 1}_img{img_idx + 1}.{image_ext}"
                    img_path = os.path.join(output_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    images.append(
                        {
                            "page": page_num + 1,
                            "path": img_path,
                            "filename": img_filename,
                            "width": width,
                            "height": height,
                            "size_bytes": len(image_bytes),
                            "ext": image_ext,
                        }
                    )
                except Exception:
                    continue
        doc.close()
    except Exception:
        pass
    return images


def extract_images_from_docx(path: str, output_dir: str = "") -> list:
    """Extract embedded images from DOCX files.

    Returns a list of dicts: [{path, filename, size_bytes, content_type}]
    """
    images = []
    try:
        from docx import Document

        doc = Document(path)
        if not output_dir:
            output_dir = os.path.join(os.path.dirname(path), "extracted_images")
        os.makedirs(output_dir, exist_ok=True)

        for idx, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type or ""
                    ext = content_type.split("/")[-1] if "/" in content_type else "png"
                    ext = ext.replace("jpeg", "jpg")

                    if len(image_bytes) < 10_000:
                        continue

                    img_filename = f"docx_img{idx + 1}.{ext}"
                    img_path = os.path.join(output_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    images.append(
                        {
                            "path": img_path,
                            "filename": img_filename,
                            "size_bytes": len(image_bytes),
                            "content_type": content_type,
                        }
                    )
                except Exception:
                    continue
    except Exception:
        pass
    return images


def extract_images_from_pptx(path: str, output_dir: str = "") -> list:
    """Extract embedded images from PPTX slides.

    Returns a list of dicts: [{slide, path, filename, size_bytes}]
    """
    images = []
    try:
        from pptx import Presentation

        prs = Presentation(path)
        if not output_dir:
            output_dir = os.path.join(os.path.dirname(path), "extracted_images")
        os.makedirs(output_dir, exist_ok=True)

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = (
                            image.content_type.split("/")[-1]
                            if image.content_type
                            else "png"
                        )
                        ext = ext.replace("jpeg", "jpg")

                        if len(image_bytes) < 10_000:
                            continue

                        img_filename = f"slide{slide_idx + 1}_img{shape_idx + 1}.{ext}"
                        img_path = os.path.join(output_dir, img_filename)
                        with open(img_path, "wb") as f:
                            f.write(image_bytes)

                        images.append(
                            {
                                "slide": slide_idx + 1,
                                "path": img_path,
                                "filename": img_filename,
                                "size_bytes": len(image_bytes),
                            }
                        )
                    except Exception:
                        continue
    except Exception:
        pass
    return images


def extract_all_images(file_path: str, output_dir: str = "") -> list:
    """Extract images from any supported file type."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_images_from_pdf(file_path, output_dir)
    elif ext in (".docx", ".doc"):
        return extract_images_from_docx(file_path, output_dir)
    elif ext in (".pptx", ".ppt"):
        return extract_images_from_pptx(file_path, output_dir)
    return []


def _extract_docx(path: str) -> str:
    try:
        from docx import Document

        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def _extract_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(path, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except Exception:
        return ""


def _extract_text(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception:
        return ""


async def ingest_document(
    doc_id: str,
    file_path: str,
    collection: str,
    metadata: Optional[Dict] = None,
    chunk_size: int = 500,
    chunk_overlap: int = 50,
) -> Dict[str, Any]:
    """Full ingestion pipeline: extract → chunk → embed → store."""
    meta = metadata or {}

    # Extract text
    text = extract_text_from_file(file_path)
    if not text.strip():
        return {"status": "error", "message": "No text content extracted", "chunks": 0}

    # Chunk
    chunks = embedding_service.chunk_text(text, chunk_size, chunk_overlap)

    # Embed and store each chunk
    for i, chunk in enumerate(chunks):
        embedding = embedding_service.embed(chunk)
        rag_pipeline.add_embedding(
            doc_id=doc_id,
            chunk_text=chunk,
            embedding=embedding,
            metadata={
                "collection": collection,
                "chunk_index": i,
                "filename": meta.get("filename", os.path.basename(file_path)),
                "outcome": meta.get("outcome"),
                "product": meta.get("product"),
                "industry": meta.get("industry"),
            },
        )

    return {
        "status": "success",
        "chunks": len(chunks),
        "text_length": len(text),
        "doc_id": doc_id,
    }


def check_duplicate(filename: str, existing_docs: List[Dict]) -> Optional[Dict]:
    """Check if a document with the same filename already exists."""
    for doc in existing_docs:
        if doc.get("filename") == filename:
            return doc
    return None
